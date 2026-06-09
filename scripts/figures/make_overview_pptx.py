"""Figure 1 (overview) — detailed 2x2 architecture for contributions C1-C4.

Editable PowerPoint deliverable. A single full-width figure with four detailed
panels arranged 2x2 as the paper's arc:
    PROBLEM (C1)  ->  DIAGNOSE (C2)
        |                 |
    SELECT  (C3)  <-  IMPROVE  (C4)   (reading arc C1->C2->C3->C4)
Each panel is a real mini-architecture (boxes + arrows + the headline number),
so each Cn reads on its own and the body prose can shrink drastically.

Real logos: official GitHub org avatars under icons/logos/norm/.
Output: paper/figures/fig1_overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
LOGO = HERE / "icons" / "logos" / "norm"
OUT = HERE.parent.parent / "paper" / "figures" / "fig1_overview.pptx"

FONT = "Calibri"
INK = "1F3247"
GREY_TX = "55636E"
GREY_LN = "C2CBD2"
ARROW = "8A98A4"

# arc colours
COL = {
    "C1": ("C62828", "FBE9EA", "FDF4F4"),   # PROBLEM  red   (header, fill, soft)
    "C2": ("1565C0", "E7EFFB", "F3F7FD"),   # DIAGNOSE blue
    "C3": ("2E7D32", "E6F3E9", "F2F9F3"),   # SELECT   green
    "C4": ("E65100", "FCEEE0", "FDF6EF"),   # IMPROVE  amber
}


def C(h):
    return RGBColor.from_string(h)


def rrect(slide, x, y, w, h, fill, line=None, lw=1.0, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def txt(shape, paras, anchor=MSO_ANCHOR.MIDDLE, m=0.04, space_after=1.5):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(m)
    tf.margin_top = tf.margin_bottom = Inches(0.015)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.CENTER)
        p.space_after = Pt(para.get("space_after", space_after))
        p.space_before = Pt(0)
        if para.get("line_spacing"):
            p.line_spacing = para["line_spacing"]
        for (t, s, b, c) in para["runs"]:
            r = p.add_run(); r.text = t
            r.font.size = Pt(s); r.font.bold = b
            r.font.name = FONT; r.font.color.rgb = C(c)
    return tf


def box(slide, x, y, w, h, lines, fill="FFFFFF", line=GREY_LN, lw=1.0, radius=0.12,
        anchor=MSO_ANCHOR.MIDDLE):
    sp = rrect(slide, x, y, w, h, fill, line=line, lw=lw, radius=radius)
    txt(sp, lines, anchor=anchor)
    return sp


def tbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.MIDDLE):
    sp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txt(sp, lines, anchor=anchor)
    return sp


def arrow(slide, x1, y1, x2, y2, color=ARROW, w=1.5):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = C(color); cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def run(t, s, b, c):
    return (t, s, b, c)


def line(runs, align=PP_ALIGN.CENTER, sa=1.5, ls=None):
    d = {"runs": runs, "align": align, "space_after": sa}
    if ls:
        d["line_spacing"] = ls
    return d


def panel_header(slide, x, y, w, key, verb, title):
    hcol = COL[key][0]
    hd = rrect(slide, x, y, w, 0.46, hcol, radius=0.12)
    txt(hd, [
        line([run(verb + "   ", 8.5, True, "FFFFFF"), run(title, 11.5, True, "FFFFFF")]),
    ])


def main():
    prs = Presentation()
    W, H = 9.6, 6.7
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title banner
    ban = rrect(slide, 0.2, 0.12, W - 0.4, 0.40, "37474F", radius=0.16)
    txt(ban, [line([
        run("RADP — choose the parser by ", 11.5, True, "FFFFFF"),
        run("retrieval", 11.5, True, "FFD54F"),
        run(", not by appearance; train only where it helps", 11.5, True, "FFFFFF"),
    ])])

    gx, gy = 0.2, 0.66
    pw, ph = 4.50, 2.86
    gap = 0.30
    pos = {
        "C1": (gx, gy),
        "C2": (gx + pw + gap, gy),
        "C3": (gx, gy + ph + gap),
        "C4": (gx + pw + gap, gy + ph + gap),
    }

    # panel frames + headers
    titles = {
        "C1": ("PROBLEM", "C1 · The disconnect"),
        "C2": ("DIAGNOSE", "C2 · Coverage diagnostic"),
        "C3": ("SELECT", "C3 · RCPS protocol"),
        "C4": ("IMPROVE", "C4 · Bounded training"),
    }
    for k, (px, py) in pos.items():
        rrect(slide, px, py, pw, ph, COL[k][2], line=COL[k][0], lw=1.3, radius=0.04)
        panel_header(slide, px, py, pw, k, *titles[k])

    def body(k):
        px, py = pos[k]
        return px + 0.14, py + 0.46 + 0.10, pw - 0.28, ph - 0.46 - 0.22

    # ---------------- C1 : disconnect ----------------
    bx, by, bw, bh = body("C1")
    box(slide, bx + bw / 2 - 0.62, by, 1.24, 0.34,
        [line([run("Parser output", 9, True, INK)])])
    # two opposing axes
    cwid = (bw - 0.18) / 2
    lx, rx = bx, bx + cwid + 0.18
    ay = by + 0.34
    box(slide, lx, ay + 0.30, cwid, 0.78, [
        line([run("Appearance", 8.5, True, "8a8f96")]),
        line([run("BC · edit-dist", 7.5, False, GREY_TX)], sa=2),
        line([run("MinerU ranks #1", 8, True, "C62828")]),
    ], fill="F4F5F6", anchor=MSO_ANCHOR.MIDDLE)
    box(slide, rx, ay + 0.30, cwid, 0.78, [
        line([run("Retrieval", 8.5, True, "2E7D32")]),
        line([run("RCPS", 7.5, False, GREY_TX)], sa=2),
        line([run("MinerU ranks LAST", 8, True, "2E7D32")]),
    ], fill="EAF4EC")
    arrow(slide, bx + bw / 2 - 0.30, ay, lx + cwid / 2, ay + 0.30)
    arrow(slide, bx + bw / 2 + 0.30, ay, rx + cwid / 2, ay + 0.30)
    box(slide, bx, by + bh - 0.40, bw, 0.38, [
        line([run("BC ↔ RCPS  r = −0.81", 9, True, "C62828"),
              run("   ·   parser choice → Hit@1 ", 8.5, False, INK),
              run("2.8×", 9.5, True, "C62828"),
              run(" (0.20→0.55)", 8.5, False, INK)]),
    ], fill="FBE9EA", line=None, radius=0.18)

    # ---------------- C2 : coverage decision tree ----------------
    bx, by, bw, bh = body("C2")
    box(slide, bx + bw / 2 - 0.72, by, 1.44, 0.32,
        [line([run("gold answer span", 8.5, True, INK)])])
    q1y = by + 0.42
    box(slide, bx + 0.30, q1y, bw - 0.60, 0.30,
        [line([run("in parser output?", 8.5, True, "1565C0")])], fill="E7EFFB")
    # NO -> absent
    box(slide, bx + bw - 1.62, q1y + 0.40, 1.62, 0.34,
        [line([run("ABSENT  20.2%", 8.5, True, "FFFFFF"), run("  parser fault", 7.5, False, "FFFFFF")])],
        fill="C62828", line=None)
    q2y = q1y + 0.40
    box(slide, bx, q2y, 1.86, 0.30,
        [line([run("fits one chunk?", 8.5, True, "1565C0")])], fill="E7EFFB")
    box(slide, bx, q2y + 0.40, 1.86, 0.34,
        [line([run("SPLIT  ≤2.3%", 8.5, True, "FFFFFF"), run("  chunker fault", 7.5, False, "FFFFFF")])],
        fill="E68A00", line=None)
    box(slide, bx + bw - 1.62, q2y + 0.40, 1.62, 0.34,
        [line([run("COVERED", 8.5, True, "FFFFFF")])], fill="2E7D32", line=None)
    arrow(slide, bx + bw / 2, by + 0.32, bx + bw / 2, q1y)
    arrow(slide, bx + bw - 0.55, q1y + 0.30, bx + bw - 0.80, q1y + 0.40, color="C62828")
    arrow(slide, bx + 0.6, q1y + 0.30, bx + 0.6, q2y, color="1565C0")
    arrow(slide, bx + 0.55, q2y + 0.30, bx + 0.55, q2y + 0.40, color="E68A00")
    arrow(slide, bx + 1.5, q2y + 0.30, bx + bw - 1.0, q2y + 0.40, color="2E7D32")
    tbox(slide, bx, by + bh - 0.26, bw, 0.24,
         [line([run("absent constant across 8 chunkers → fix the parser, not the chunker", 7.5, False, GREY_TX)])])

    # ---------------- C3 : RCPS protocol ----------------
    bx, by, bw, bh = body("C3")
    box(slide, bx, by, bw, 0.30,
        [line([run("held-out Q–A probe  +  parsed corpus", 8.5, True, INK)])])
    box(slide, bx + bw / 2 - 0.9, by + 0.40, 1.8, 0.28,
        [line([run("chunk → retrieve (MRR)", 8, True, "2E7D32")])], fill="EAF4EC")
    cy = by + 0.78
    cw3 = (bw - 0.16) / 3
    for i, t in enumerate(["✓ extrinsic\n(probe)", "✓ retriever-avg\n×3", "✓ format-\ninvariant"]):
        a, b2 = t.split("\n")
        box(slide, bx + i * (cw3 + 0.08), cy, cw3, 0.52,
            [line([run(a, 7.8, True, "2E7D32")], sa=1), line([run(b2, 7.2, False, GREY_TX)])],
            fill="F2F9F3")
    box(slide, bx + 0.2, by + bh - 0.62, bw - 0.4, 0.34,
        [line([run("RCPS ranking → pick parser + chunker  (no training)", 8.3, True, "2E7D32")])],
        fill="E6F3E9", line=None)
    tbox(slide, bx, by + bh - 0.24, bw, 0.22,
         [line([run("retriever-avg flips the top parser a single embedder gets wrong (τ = 0.87)", 7.3, False, GREY_TX)])])
    arrow(slide, bx + bw / 2, by + 0.30, bx + bw / 2, by + 0.40, color="2E7D32")
    arrow(slide, bx + bw / 2, by + 0.68, bx + bw / 2, cy, color="2E7D32")

    # ---------------- C4 : training pipeline ----------------
    bx, by, bw, bh = body("C4")
    sw = (bw - 3 * 0.22) / 4
    yy = by + 0.06
    steps = ["Prod\nparser", "sample K\nparses", "rank\ncandidates", "LoRA-toggle\nDPO"]
    for i, t in enumerate(steps):
        a, b2 = t.split("\n")
        box(slide, bx + i * (sw + 0.22), yy, sw, 0.56,
            [line([run(a, 7.8, True, INK)], sa=1), line([run(b2, 7.2, False, GREY_TX)])],
            fill="FDF6EF" if i else "FFFFFF")
        if i:
            arrow(slide, bx + i * (sw + 0.22) - 0.21, yy + 0.28, bx + i * (sw + 0.22) + 0.01, yy + 0.28,
                  color="E65100")
    # ranking routes annotation
    tbox(slide, bx, yy + 0.60, bw, 0.40, [
        line([run("rank by: ", 7.8, True, INK),
              run("edit-dist → GT", 8, True, "2E7D32"),
              run("  (RADP-Distill ✓)", 7.8, True, "2E7D32")]),
        line([run("vs  page-RCPS  (RADP-DPO) — same result", 7.8, False, GREY_TX)], sa=1),
    ], anchor=MSO_ANCHOR.TOP)
    box(slide, bx, by + bh - 0.42, bw, 0.40, [
        line([run("retrieval reward unnecessary", 8.5, True, "E65100"),
              run("  →  distil to clean GT text", 8.3, False, INK)]),
        line([run("+1.22 pp Hit@5 (OHR-Bench)", 8.5, True, "E65100")], sa=1),
    ], fill="FCEEE0", line=None, radius=0.16)

    # arc arrows: within-row (C1->C2, C3->C4) + a centred down arrow joining the
    # "diagnose" row to the "act" row (cleaner than a diagonal C2->C3 wrap)
    arrow(slide, pos["C1"][0] + pw + 0.02, gy + ph / 2, pos["C2"][0] - 0.02, gy + ph / 2, color="6b7780", w=2.0)
    arrow(slide, pos["C3"][0] + pw + 0.02, gy + ph + gap + ph / 2, pos["C4"][0] - 0.02,
          gy + ph + gap + ph / 2, color="6b7780", w=2.0)
    arrow(slide, W / 2, gy + ph + 0.005, W / 2, gy + ph + gap - 0.005, color="6b7780", w=2.0)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("saved ->", OUT)


if __name__ == "__main__":
    main()
