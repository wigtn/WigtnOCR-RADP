"""Figure 1 (overview) as an EDITABLE PowerPoint, OHR-Bench-Fig2 density.

Full-width, three phase columns, each packed with REAL artifacts (page image,
the garbled-vs-clean parse contrast, the actual result charts, brand logos) so
it reads as densely as OHR-Bench Fig.2 while staying organised:

  C1 Parsing-Retrieval Disconnect  ->  C2/C3 Diagnose + Select  ->  C4 Bounded Training

No heavy bottom bar -- a thin rule + one bold takeaway line.
Korean is set as latin+EA font (Noto Sans CJK KR here; swap in PowerPoint).
NOTE: brand logos under scripts/figures/icons/logos/ are NOT committed (licensing);
download from each project's official site before running. See README "figure logos".
Output: paper/figures/fig_overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = Path("paper/figures/fig_overview.pptx")
FIG = Path("paper/figures")
LOGO = Path("scripts/figures/icons/logos")
PAGE = Path("data/KoGovDoc-Bench/images/documents/kogov_003/page_0533.png")

KFONT, CF = "Noto Sans CJK KR", "Calibri"
WHITE = RGBColor(0xFF, 0xFF, 0xFF); DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x55, 0x55, 0x55); LINE = RGBColor(0x90, 0x90, 0x90)
BLUE = RGBColor(0x1F, 0x5F, 0xA8); BLUEF = RGBColor(0xEC, 0xF2, 0xFB)
AMBER = RGBColor(0x9A, 0x5A, 0x00); AMBERF = RGBColor(0xFC, 0xF3, 0xE6)
GREEN = RGBColor(0x1A, 0x6B, 0x1A); GREENF = RGBColor(0xEC, 0xF6, 0xEC)
RED = RGBColor(0x8F, 0x1B, 0x1B); REDF = RGBColor(0xF9, 0xEC, 0xEC)
GBAR = RGBColor(0x4A, 0x9A, 0x4A); ABAR = RGBColor(0xD9, 0x8A, 0x2B)
RBAR = RGBColor(0xBE, 0x3A, 0x3A)


def _ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for t in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(t))
        if el is None:
            el = rPr.makeelement(qn(t), {}); rPr.append(el)
        el.set("typeface", name)


def R(txt, size, color, bold=False, italic=False, font=CF):
    return (txt, size, color, bold, italic, font)


def tbox(s, x, y, w, h, paras, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         lead=1.1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = lead
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; _ea(r, font)
    return tb


def rrect(s, x, y, w, h, fill, line, lw=1.25, radius=0.04):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.adjustments[0] = radius
    sp.fill.background() if fill is None else (sp.fill.solid(),
                                               setattr(sp.fill.fore_color, "rgb", fill))
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def chevron(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def bar(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(max(w, 0.03)), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.color.rgb = WHITE; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def pic_fit(s, path, x, y, w, h):
    """Place an image scaled to fit (w,h), centred in that box."""
    iw, ih = Image.open(path).size
    sc = min(w / iw, h / ih)
    dw, dh = iw * sc, ih * sc
    return s.shapes.add_picture(str(path), Inches(x + (w - dw) / 2),
                                Inches(y + (h - dh) / 2), height=Inches(dh))


def logo(s, name, x, y, h):
    return s.shapes.add_picture(str(LOGO / name), Inches(x), Inches(y),
                                height=Inches(h))


def subzone(s, x, y, w, h, title, color, fill, lw=1.25, tsize=10.5):
    rrect(s, x, y, w, h, WHITE, color, lw=lw)
    rrect(s, x, y, w, 0.34, fill, None, radius=0.10)
    tbox(s, x + 0.05, y, w - 0.1, 0.34, [[R(title, tsize, color, bold=True)]],
         anchor=MSO_ANCHOR.MIDDLE)


def cell(tbl, r, c, text, *, size=9, color=DARK, bold=False, fill=None,
         font=CF, align=PP_ALIGN.CENTER):
    cl = tbl.cell(r, c)
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_left = cl.margin_right = Pt(3)
    cl.margin_top = cl.margin_bottom = Pt(1)
    if fill is not None:
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
    p = cl.text_frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; _ea(run, font)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.3); prs.slide_height = Inches(5.7)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    M, GAP = 0.22, 0.26
    TOT = 13.3 - 2 * M - 2 * GAP
    w1, w2, w3 = TOT * 0.95 / 3, TOT * 1.20 / 3, TOT * 0.85 / 3   # C3 column widest
    x1 = M; x2 = M + w1 + GAP; x3 = M + w1 + GAP + w2 + GAP

    # ---- column header bars (C3 = central contribution -> widest, biggest type) ----
    for x, w, txt, c, fs in [(x1, w1, "C1 · Parsing–Retrieval Disconnect", BLUE, 11.5),
                             (x2, w2, "C3 · RCPS Selection  (+ C2 Diagnose)", AMBER, 13),
                             (x3, w3, "C4 · Bounded Parser Training", GREEN, 11.5)]:
        rrect(s, x, 0.2, w, 0.46, c, None, radius=0.12)
        tbox(s, x, 0.2, w, 0.46, [[R(txt, fs, WHITE, bold=True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    top, BOT = 0.82, 5.56   # all three columns share this bottom edge
    # ================= COLUMN 1 =================
    # 1a candidate parsers (logo + name under each; no stray "+GOT" text)
    h1a = 1.18
    subzone(s, x1, top, w1, h1a, "Candidate parsers", BLUE, BLUEF)
    for i, (nm, lab) in enumerate([("qwen.png", "Qwen3-VL"),
            ("mineru.png", "MinerU"), ("marker_datalab.png", "Marker"),
            ("paddle.png", "PaddleOCR")]):
        dx = 0.45 + i * 0.9
        logo(s, nm, x1 + dx, top + 0.40, 0.42)
        tbox(s, x1 + dx - 0.27, top + 0.86, 0.96, 0.24,
             [[R(lab, 7.5, GREY)]], align=PP_ALIGN.CENTER)

    # 1b cleaner-looking != better retrieval (clean comparison table; no
    # page screenshot, no raw multilingual parse text)
    y1b = top + h1a + 0.16; h1b = BOT - y1b
    subzone(s, x1, y1b, w1, h1b, "Cleaner-looking ≠ better retrieval", BLUE, BLUEF)
    tw = w1 - 0.4
    gt = s.shapes.add_table(3, 3, Inches(x1 + 0.2), Inches(y1b + 0.58),
                            Inches(tw), Inches(1.6)).table
    gt.first_row = False; gt.horz_banding = False
    wts = [1.05, 1.45, 0.95]
    for i, wv in enumerate(wts):
        gt.columns[i].width = Inches(tw * wv / sum(wts))
    GH = RGBColor(0xEC, 0xEC, 0xEC)
    rows = [
        [("parser", DARK, True, GH), ("Boundary Clarity", DARK, True, GH),
         ("Hit@1", DARK, True, GH)],
        [("MinerU", DARK, True, REDF), ("0.72  (highest)", DARK, False, REDF),
         ("0.20", RED, True, REDF)],
        [("Prod (ours)", DARK, True, GREENF), ("0.61", DARK, False, GREENF),
         ("0.55", GREEN, True, GREENF)],
    ]
    for r, row in enumerate(rows):
        for c, (txt, col, bd, fl) in enumerate(row):
            cell(gt, r, c, txt, size=(8.5 if r == 0 else 9.5), color=col,
                 bold=bd, fill=fl)
    tbox(s, x1 + 0.14, y1b + h1b - 1.05, w1 - 0.28, 0.95,
         [[R("parser choice alone moves Hit@1 by ", 10.5, DARK),
           R("+35.1pp (0.20→0.55; 2.8×)", 12, BLUE, bold=True)],
          [R("appearance metrics measure formatting, not content", 9, DARK)],
          [R("(anti-correlation  r = −0.81,  n = 5)", 8.5, GREY, italic=True)]],
         align=PP_ALIGN.CENTER, lead=1.4)

    # ================= COLUMN 2 (C3 central, C2 support) =================
    h2a = 2.70   # C3 = central contribution -> the larger, top box
    subzone(s, x2, top, w2, h2a, "C3 · RCPS selection protocol (no training)", AMBER, AMBERF, lw=2.0, tsize=11)
    tbox(s, x2 + 0.12, top + 0.42, w2 - 0.24, 0.4,
         [[R("held-out Q–A probe, scored over 3 retrievers:", 9.5, DARK)]])
    for nm, dx, lab in [("bge_baai.png", 1.0, "BGE-M3"),
                        ("me5_ms.png", 2.1, "mE5"), ("qwen.png", 3.2, "Qwen3-Emb")]:
        logo(s, nm, x2 + dx, top + 0.82, 0.48)
        tbox(s, x2 + dx - 0.21, top + 1.32, 0.9, 0.22,
             [[R(lab, 8, GREY, italic=True)]], align=PP_ALIGN.CENTER)
    tbox(s, x2 + 0.12, top + 1.64, w2 - 0.24, 0.26,
         [[R("extrinsic · retriever-averaged · format-normalised", 9.2, DARK, bold=True)]],
         align=PP_ALIGN.CENTER)
    tbox(s, x2 + 0.12, top + 1.92, w2 - 0.24, 0.24,
         [[R("no training · no manual relevance labels", 8.8, GREY, italic=True)]],
         align=PP_ALIGN.CENTER)
    tbox(s, x2 + 0.12, top + h2a - 0.52, w2 - 0.24, 0.46,
         [[R("→ rank parsers & chunkers by retrieval", 10.5, AMBER, bold=True)],
          [R("picks Prod (RCPS 0.583) over MinerU (0.212)", 9, DARK)]],
         align=PP_ALIGN.CENTER, lead=1.25)

    y2b = top + h2a + 0.16; h2b = BOT - y2b   # C2 = compact support box
    subzone(s, x2, y2b, w2, h2b, "C2 · Coverage diagnostic (no retriever)", AMBER, AMBERF)
    bx, bw, by, bh = x2 + 0.25, w2 - 0.5, y2b + 0.46, 0.5
    cx = bx
    for frac, col in [(0.775, GBAR), (0.023, ABAR), (0.202, RBAR)]:
        bar(s, cx, by, bw * frac, bh, col); cx += bw * frac
    tbox(s, bx, by, bw * 0.775, bh, [[R("77.5%", 9.5, WHITE, bold=True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, bx + bw * 0.798, by, bw * 0.202, bh,
         [[R("20.2%", 9.5, WHITE, bold=True)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, bx, by + bh + 0.03, bw * 0.775, 0.22,
         [[R("covered", 8.5, GBAR, bold=True)]], align=PP_ALIGN.CENTER)
    tbox(s, bx + bw * 0.70, by + bh + 0.03, bw * 0.30, 0.22,
         [[R("absent", 8.5, RBAR, bold=True)]], align=PP_ALIGN.CENTER)
    tbox(s, x2 + 0.14, by + bh + 0.28, w2 - 0.28, 0.5,
         [[R("absent = ", 9, DARK), R("parser-side failure", 9, RED, bold=True),
           R(", set before chunking", 9, DARK)],
          [R("no chunker can recover it → fix the parser (C4)", 8.8,
             GREY, italic=True)]], align=PP_ALIGN.CENTER, lead=1.25)

    # ================= COLUMN 3 (C4) =================
    h3a = 1.95
    subzone(s, x3, top, w3, h3a, "C4 · Parser-side training", GREEN, GREENF)
    tbox(s, x3 + 0.14, top + 0.44, w3 - 0.28, 0.4,
         [[R("best-of-K parses from Prod, ranked two ways:", 9.5, DARK)]])
    rrect(s, x3 + 0.2, top + 0.86, w3 - 0.4, 0.46, GREENF, GREEN, 1.0)
    tbox(s, x3 + 0.22, top + 0.86, w3 - 0.44, 0.46,
         [[R("RADP-Distill — edit-distance → ref", 9.5, GREEN, bold=True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, x3, top + 1.34, w3, 0.22, [[R("≈", 13, DARK, bold=True)]],
         align=PP_ALIGN.CENTER)
    rrect(s, x3 + 0.2, top + 1.42, w3 - 0.4, 0.42, AMBERF, AMBER, 1.0)
    tbox(s, x3 + 0.22, top + 1.42, w3 - 0.44, 0.42,
         [[R("RADP-DPO — page-local RCPS", 9.5, AMBER, bold=True)]],
         anchor=MSO_ANCHOR.MIDDLE)

    y3b = top + h3a + 0.16; h3b = BOT - y3b
    subzone(s, x3, y3b, w3, h3b, "Result", GREEN, GREENF)
    tbox(s, x3 + 0.14, y3b + 0.42, w3 - 0.28, 0.3,
         [[R("Hit@5 gain over the untuned parser:", 9.5, DARK)]])
    blx, sc = x3 + 1.25, 1.30 / 1.22
    for ry, nm, val, col in [(y3b + 0.95, "RADP-Distill", 1.22, GBAR),
                             (y3b + 1.42, "RADP-DPO", 0.85, ABAR)]:
        tbox(s, x3 + 0.14, ry - 0.03, 1.05, 0.32,
             [[R(nm, 9, col, bold=True)]], anchor=MSO_ANCHOR.MIDDLE)
        bar(s, blx, ry, val * sc, 0.3, col)
        tbox(s, blx + val * sc + 0.05, ry - 0.03, 0.85, 0.32,
             [[R(f"+{val:.2f} pp", 9.5, DARK, bold=True)]],
             anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, x3 + 0.14, y3b + h3b - 0.66, w3 - 0.28, 0.6,
         [[R("substantially overlapping CIs", 9.5, DARK, bold=True)],
          [R("no evidence reward helps; lever = fidelity distillation", 9,
             GREEN, bold=True)]], align=PP_ALIGN.CENTER, lead=1.28)

    # ---- flow chevrons between columns ----
    chevron(s, x1 + w1 + 0.01, 2.75, GAP - 0.02, 0.6, AMBER)
    chevron(s, x2 + w2 + 0.01, 2.75, GAP - 0.02, 0.6, GREEN)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"saved -> {OUT}")


if __name__ == "__main__":
    main()
